#!/usr/bin/env python3
"""
Builds DSCI 504 NorthBay Pantry presentation as a .pptx file.
Upload the output file to Google Drive → Google Slides will open it directly.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colors
NAVY  = RGBColor(0x1a, 0x23, 0x32)
CYAN  = RGBColor(0x00, 0xD4, 0xFF)
LGRAY = RGBColor(0xE8, 0xEA, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xFF, 0xCC, 0x33)
DKBLUE = RGBColor(0x0F, 0x34, 0x60)

SW = Inches(13.33)
SH = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def set_bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                color=LGRAY, size=18, bold=False, align=PP_ALIGN.LEFT,
                line_spacing=None):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align

    # Handle multi-line text — split on newline
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            run = p.add_run()
        else:
            p = tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
        run.text = line
        run.font.color.rgb = color
        run.font.size = Pt(size)
        run.font.bold = bold

    if line_spacing:
        from pptx.util import Pt as P
        from pptx.oxml.ns import qn
        from lxml import etree
        for para in tf.paragraphs:
            pPr = para._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
            spcPct.set('val', str(int(line_spacing * 1000)))

    return txb


def title_box(slide, text, size=34):
    add_textbox(slide, text,
                left=Inches(0.5), top=Inches(0.25),
                width=Inches(12.33), height=Inches(1.0),
                color=CYAN, size=size, bold=True)


def body_box(slide, text, size=15, top=Inches(1.4), height=Inches(5.7)):
    add_textbox(slide, text,
                left=Inches(0.5), top=top,
                width=Inches(12.33), height=height,
                color=LGRAY, size=size)


# ── Slide definitions ─────────────────────────────────────────────────────────

SLIDES = [
    {
        'type': 'title',
        'main': 'NorthBay Pantry\nEnterprise Data Platform',
        'sub': 'Designing an Enterprise Data Platform from Scratch',
        'byline': 'DSCI 504  ·  UBA MADS  ·  Ahmad Var  ·  May 2026',
    },
    {
        'title': 'The Problem — $1.4B in Revenue, Zero Data Infrastructure',
        'body': """160 stores across US + UK  ·  E-commerce  ·  Pantry Box meal-kit  ·  PantryPlus loyalty

10 data sources:
  • Some streaming at 80,000 events/minute
  • Some nightly CSV files with garbled headers and encoding issues
  • Some WAV audio files from 250 hours/day of customer support calls

Finance can't see revenue until spreadsheets are emailed each morning
Customer Success doesn't know a caller's order history when they pick up the phone
Data Science can't train models — no consistent, centralized data access

→  This is a greenfield build at $1.4B enterprise scale. No data warehouse. No data team.""",
        'body_size': 17,
    },
    {
        'title': 'Design Philosophy — Match Tool to Workload',
        'body': """A single database cannot serve real-time streaming, sub-second dashboards, and 5-year model training simultaneously.

Three questions drove every design decision:

  1.  What is the shape of the data?       Structured / semi-structured / binary (audio, HTML)
  2.  What is the access pattern?          Streaming / batch / sub-second online lookup
  3.  Who owns the data and who can see it?

Engineering decisions are principled — not arbitrary.
Every technology choice traces back to one of these three questions.""",
        'body_size': 17,
    },
    {
        'title': 'D1 — Reference Architecture: Bronze → Silver → Gold',
        'body': """  10 Sources  →  Ingestion Layer  →  Bronze (GCS)  →  Silver (BigQuery)  →  Gold (BigQuery)  →  Consumers
                                           ↑                  ↑                    ↑
                                      Raw, immutable     Cleaned, validated    Business-ready
                                      never modified     bad rows quarantined   role-gated

Bronze (GCS)          Raw files, forever. The recovery layer — any downstream bug is fixed by reprocessing from here.
Silver (BigQuery)     dbt cleans, validates, deduplicates. Entity resolution produces one customer_uuid per person.
Gold (BigQuery)       Finance, Marketing, and Data Science query here. Access enforced by role.
Feature Store         Feast + Redis — ML features served in <5ms to meet the Customer Success 2-second SLA.

Orchestration:  Apache Airflow (batch pipelines)  +  Apache Kafka (streaming events)""",
        'body_size': 14,
    },
    {
        'title': 'D2 — Ingestion: 10 Sources, 3 Patterns',
        'body': """STREAMING (real-time)
  Sources:  OLTP PostgreSQL (CDC), Mobile App Events, MongoDB Change Streams
  Tool:     Debezium → Kafka → Kafka Connect → Bronze

BATCH (scheduled)
  Sources:  Store POS CSV/SFTP (160 stores), Salesforce, PantryPlus Loyalty, 4 Marketing Platforms, Supplier EDI
  Tool:     Apache Airflow — one DAG per source

EVENT-DRIVEN (continuous)
  Sources:  Customer Support Call WAV files (~250 hrs/day)
  Tool:     S3 event → SQS → Lambda → AWS Transcribe → BigQuery

KEY DESIGN DECISIONS
  Fault isolation on POS     One Airflow task per store. One broken CSV → quarantine + alert. 159 stores continue unaffected.
  Idempotency everywhere     Bronze partitioned by source/YYYY/MM/DD/run_id — re-runs overwrite the same partition, never double-count.
  Rate limiting              PantryPlus API capped at 60 req/min via token-bucket limiter built into the Airflow operator.""",
        'body_size': 14,
    },
    {
        'title': 'D3 — Storage Architecture: Right Tool for the Right Workload',
        'body': """GCS / S3               Raw files, WAV audio, scraped HTML
                       Cheapest $/GB · no schema enforcement · handles any binary format

Apache Kafka           In-flight events from OLTP, Mobile App, MongoDB
                       Decouples producers from consumers · replay on failure · Schema Registry enforces Avro contracts

BigQuery (Silver/Gold) Cleaned tables, dashboards, all analytical queries
                       Columnar · TB-scale · sub-second aggregations · native dbt integration

Feast + Redis          Online ML feature store — <5ms feature lookup for Customer Success
                       BigQuery's ~1–3 sec query latency violates the 2-second SLA. Redis solves it.

Feast + BigQuery       Offline ML feature store — full history for training churn, demand forecast, recommender

"Why not one system?"
  BigQuery can't store WAV files.  S3 can't answer SELECT SUM(revenue) GROUP BY store in milliseconds.
  Redis exists for one specific reason: to serve a single SLA that BigQuery cannot meet.""",
        'body_size': 14,
    },
    {
        'title': 'D4 — Curation: The Customer Identity Problem',
        'body': """Same person — 6 different records across 6 systems:
  "ahmad.var@gmail.com"   in OLTP PostgreSQL
  "AHMAD VAR" + phone     in Salesforce
  "ahmadvar@gmail.com"    in PantryPlus loyalty (typo variant)
  Device ID only          in Mobile App
  Hashed email            in Marketing platforms
  Username (pseudonymous) in Customer Reviews

Entity Resolution Pipeline:
  1  Deterministic pass     Exact match on normalized email → lower().strip() → same customer, guaranteed
  2  Probabilistic pass     Splink scoring: name + phone Soundex + address zip ≥ 0.85 → merge candidate
  3  Golden record          One customer_uuid (UUID v4) written to golden_customer_master
  4  Survivorship rule      Most recently updated non-null field wins · OLTP email is canonical

Also:
  golden_product_master     MongoDB is system of record · flattens nested attrs · standardizes FALCPA allergen codes
  All amounts → USD         Daily ECB rate · original currency stored alongside
  All timestamps → UTC      Display timezone applied at the query/API layer only""",
        'body_size': 14,
    },
    {
        'title': 'D5 — Data Quality: Bad Data Has a Home, But Not in Production',
        'body': """Schema conformance             Bronze  ·  Great Expectations       → Quarantine file  ·  Slack #data-alerts
Null rate on required fields   Silver  ·  dbt not_null             → Block Silver load  ·  PagerDuty P2
Duplicate primary keys         Silver  ·  dbt unique               → Block load  ·  alert
Row count ±20% vs prior day    Silver  ·  dbt custom               → PagerDuty P1 if Finance table
Referential integrity (FK)     Silver  ·  dbt relationships        → Log violation  ·  do not block
PII in non-PII table           Gold    ·  GE custom expectation    → Block  + incident ticket
Finance table immutability     Gold    ·  audit trigger            → Block  +  SOX violation alert
Freshness SLA miss             Gold    ·  Airflow SLA callback     → PagerDuty P1 if 7AM Finance miss

Quarantined rows land in bq_quarantine with reason code + timestamp.
Never deleted — reviewed weekly by data stewards.

Results surfaced in DataHub data quality dashboard, visible to all owners and stewards.""",
        'body_size': 14,
    },
    {
        'title': 'D6 + D7 — Catalog and Governance',
        'body': """DataHub (open-source catalog)
  Every column tagged:  domain · sensitivity (public / internal / confidential / restricted) · pii · source_system · owner
  Source-to-Gold lineage auto-ingested from dbt — full column-level lineage without manual effort
  PII tags flow: dbt meta:{pii:true} → DataHub → BigQuery column-level policy tag enforcement

Governance — 3 Roles, Clear Accountability
  Data Owner (VP)               Accountable for accuracy and access policy for their domain
  Data Steward (Analyst)        Reviews quarantine queue · approves schema change requests
  Data Custodian (Data Eng)     Implements access controls · runs and monitors pipelines

RACI Highlight
  golden_customer_master    Owner: VP Customer Success  |  Steward: CS Analyst    |  Legal: Consulted
  finance_daily_revenue     Owner: CFO                  |  Steward: Finance        |  Legal: Informed
  stg_call_transcripts      Owner: VP Customer Success  |  Steward: CS Analyst    |  Legal: ACCOUNTABLE

  Voice recordings are the only dataset where Legal is Accountable — not just Informed.""",
        'body_size': 14,
    },
    {
        'title': 'D8 — Security and Compliance: Two Frameworks, One Platform',
        'body': """Access Controls
  BigQuery column-level policy tags:  call transcripts readable only by role:customer_success
  AES-256 at rest on all GCS + BigQuery  ·  WAV files use CMEK (customer-managed keys)  ·  TLS 1.3 in transit
  Cloud Audit Logs retained 7 years (SOX)
  SOX tables (finance_daily_revenue): append-only · no UPDATE/DELETE · change-controlled via dbt PR approval

GDPR Right-to-Erasure Runbook — target 5 business days, hard limit 30 days
  1  Look up customer_uuid in golden_customer_master using email or phone
  2  DataHub lineage API → lists every table containing customer_uuid FK — in seconds, not hours
  3  Erasure dbt macro: nullify PII columns in Silver → [ERASED-{date}]
  4  Delete WAV files + transcripts from S3/BigQuery using call UUID index
  5  Submit deletion requests to Salesforce API + PantryPlus SaaS vendor
  6  Bronze: PII stays (immutable raw) — permissible under GDPR audit exception · encrypted + access-restricted
  7  Log completion in gdpr_erasure_log · notify customer within 30-day SLA

The DataHub lineage answer to "find every row about customer X" takes seconds, not hours.""",
        'body_size': 13,
    },
    {
        'title': 'D9 — Monitoring and Operations',
        'body': """Pipeline SLAs
  Finance daily revenue     Available by 07:00 ET         P1 — page on-call immediately
  Store POS ingestion       Complete by 05:00 ET          P1 — Finance SLA depends on it
  Silver refresh            Complete by 06:00 ET          P1
  Customer 360 lookup       < 2 sec response              P2 — page within 15 min
  Loyalty / Marketing       Complete by 09:00 ET          P2

Alerting Chain
  Airflow SLA miss callback  →  PagerDuty (P1 or P2)  →  On-call data engineer  →  Slack #data-incidents

On-call: Weekly rotation among 3 engineers. Each DAG has a runbook_url tag — one click to the runbook.

Key Metrics (Airflow + Datadog)
  DAG success rate · Bronze→Silver lag (minutes) · quarantine rate per source
  Kafka consumer lag · BigQuery slot utilization + cost per DAG · Gold table freshness

Cost tagging on every job: team / source_system / environment / dag_id → monthly VP-level cost reports""",
        'body_size': 14,
    },
    {
        'title': 'D10 — Sample Artifacts: The Design Is Backed by Working Code',
        'body': """artifacts/dags/pos_ingestion_dag.py
  Per-store fault isolation via try/except — one broken file cannot stop the other 159 stores
  chardet for encoding detection · quarantine bucket + Slack alert on failure

artifacts/dbt/stg_pos_transactions.sql
  Bronze → Silver transformation: SAFE_CAST on every column · dedup via row_number()
  Incremental MERGE on natural key — guarantees idempotency on any re-run

artifacts/great_expectations/pos_suite.py
  DQ checks: nulls on required fields · unique transaction IDs · revenue in valid range
  Row count freshness check — fails if today's batch is missing entirely

artifacts/dags/alert_callback.py
  SLA miss callback: fires PagerDuty P1 payload + posts to Slack #data-incidents automatically
  Includes last-success timestamp in the alert so on-call engineer has immediate context

Each artifact demonstrates a different layer of the platform.
The per-store try/except in the DAG is the most readable illustration of fault isolation.""",
        'body_size': 15,
    },
    {
        'title': 'D11 — Cost and Risk: ~$20K–$33K/month',
        'body': """Component                      Monthly Estimate
  BigQuery (storage + queries)   $4,000 – $8,000
  GCS / S3 (Bronze + audio)      $2,000 – $3,500
  Kafka (Confluent Cloud)        $3,000 – $5,000
  Airflow (Cloud Composer)       $2,000 – $3,000
  AWS Transcribe (250 hrs/day)   $8,000 – $12,000   ← DOMINANT COST (40% of total)
  DataHub + Redis                $1,000 – $2,000
  TOTAL                          ~$20,000 – $33,500/month

Mitigation: Transcribe only escalated calls → reduce audio cost by ~60%

Top 3 Risks
  HIGH / HIGH          POS schema drift breaks Silver
                       → Schema evolution tests in CI · unknown columns quarantined on arrival

  MEDIUM / MEDIUM      Kafka consumer lag during mobile event spike
                       → Auto-scaling consumer group · 7-day topic retention for replay

  LOW / CRITICAL       GDPR erasure misses a table
                       → DataHub lineage + quarterly erasure drill + Legal audit of erasure log""",
        'body_size': 13,
    },
    {
        'title': '8 Concrete Questions — Rapid Fire',
        'body': """1   How does data flow from each source?
    D2 ingestion table — 3 patterns, 10 sources, specific tool per source

2   Where does each dataset live and why?
    D3 storage table — each system exists because no other system can do its job

3   How do you prevent double-counting on re-runs?
    Bronze overwrite partitions (source/YYYY/MM/DD/run_id) + Silver MERGE on natural key

4   How do you identify one customer across 10 sources?
    customer_uuid via deterministic email match first, then Splink probabilistic scoring (≥0.85)

5   How do you find every row about customer X in under one hour for GDPR?
    DataHub lineage API lists all tables with customer_uuid FK in seconds → erasure macro runs in minutes

6   Which datasets are most sensitive and how is access controlled?
    Support call transcripts (Restricted) + customer PII (Confidential) — BigQuery column-level policy tags

7   If Finance dashboard is missing at 7 AM — what happens?
    Airflow SLA miss callback → PagerDuty P1 → on-call engineer with runbook URL in the alert body

8   What would this cost to run?
    ~$20K–$33K/month · audio transcription is the cost lever (40% of total, reducible by ~60%)""",
        'body_size': 16,
    },
    {
        'title': 'Architecture Recap — End to End',
        'body': """  10 Sources  →  Ingestion Layer  →  Bronze (GCS)  →  Silver (BigQuery)  →  Gold (BigQuery)
        |                   |                 |                  |                   |
   Kafka path          Airflow path       Raw files          dbt models          Business tables
   (streaming)          (batch)           immutable        + Great Exp.         Finance · Mktg
   OLTP, Mobile,       POS, EDI,          forever          quarantine bad         DS queries
   MongoDB             Salesforce                           rows, never delete
        |
   Transcribe path                                                           Feature Store
   (audio events)                                                            Feast + Redis
   WAV → S3 → Lambda                                                          <5ms online
   → AWS Transcribe                                                            lookup

Key paths:
  Real-time events  →  Kafka  →  Kafka Connect  →  Bronze  →  Silver streaming model
  Daily batches     →  Airflow DAGs (one per source)  →  Bronze  →  dbt Silver  →  Gold
  Audio calls       →  S3 event trigger  →  Lambda  →  AWS Transcribe  →  BigQuery""",
        'body_size': 13,
    },
    {
        'title': 'NorthBay — From Zero Infrastructure to Enterprise-Grade Platform',
        'body': """  ✓  10 sources ingested reliably and idempotently
  ✓  One trusted customer identity across all systems
  ✓  Finance gets their 7 AM dashboard — automated runbook fires if it's ever late
  ✓  Data Science has a feature store with sub-millisecond online serving
  ✓  Legal has a GDPR runbook that completes in 5 business days, not weeks
  ✓  ~$20K–$33K/month — with a clear cost lever on audio transcription

Every design decision is auditable.
Every dataset has an owner.
Every alert has a runbook.

The platform is not just built — it's operable.""",
        'body_size': 19,
    },
]


def build():
    prs = new_prs()
    out = '/home/ahmadavar/uba-mads-projects/03-northbay-data-platform/NorthBay_DSCI504.pptx'

    for i, s in enumerate(SLIDES):
        slide = blank_slide(prs)
        set_bg(slide)

        if s.get('type') == 'title':
            # Big centered title slide
            add_textbox(slide, s['main'],
                        Inches(0.5), Inches(1.5), Inches(12.33), Inches(2.0),
                        color=CYAN, size=52, bold=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, s['sub'],
                        Inches(0.5), Inches(3.8), Inches(12.33), Inches(0.7),
                        color=WHITE, size=22, align=PP_ALIGN.CENTER)
            add_textbox(slide, s['byline'],
                        Inches(0.5), Inches(4.7), Inches(12.33), Inches(0.5),
                        color=LGRAY, size=18, align=PP_ALIGN.CENTER)
        else:
            title_box(slide, s['title'], size=min(30, 34 - max(0, len(s['title']) - 60) // 5))
            body_box(slide, s['body'], size=s.get('body_size', 15))

        # Slide number (bottom right, small)
        add_textbox(slide, str(i + 1),
                    Inches(12.5), Inches(7.1), Inches(0.4), Inches(0.3),
                    color=RGBColor(0x55, 0x66, 0x77), size=10, align=PP_ALIGN.RIGHT)

    prs.save(out)
    print(f'Saved: {out}')
    print(f'\nNext step: download this file to your laptop, then')
    print(f'upload it to Google Drive — it opens directly in Google Slides.')


if __name__ == '__main__':
    build()
